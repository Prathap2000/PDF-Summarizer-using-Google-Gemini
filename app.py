import os
import streamlit as st
import google.generativeai as genai
from dotenv import load_dotenv
import fitz  # PyMuPDF for PDF text extraction
from docx import Document  # For DOCX file handling
import io  # For handling in-memory bytes
from pptx import Presentation  # For PPTX file handling

# Load environment variables from .env file
load_dotenv()

# Configure API key for Google Gemini
genai.configure(api_key=os.getenv('API_KEY'))

# Streamlit app layout
st.title('File Summarizer using Google Gemini')

# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    text = ""
    with fitz.open(pdf_path) as pdf_file:
        for page in pdf_file:
            text += page.get_text()
    return text

# Function to extract text from DOCX
def extract_text_from_docx(docx_path):
    text = ""
    doc = Document(docx_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_path):
    text = ""
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to wrap and justify text to fit a given width
def wrap_and_justify_text(text, width):
    words = text.split()
    wrapped_lines = []
    current_line = ""

    for word in words:
        # Check if adding the next word would exceed the width
        if len(current_line) + len(word) + 1 > width:
            wrapped_lines.append(current_line.strip())  # Append the current line to the list
            current_line = word  # Start a new line with the current word
        else:
            current_line += (" " + word) if current_line else word  # Add word to the current line

    # Add any remaining text as the last line
    if current_line:
        wrapped_lines.append(current_line.strip())

    # Justify each line (except the last one)
    justified_lines = []
    for i in range(len(wrapped_lines)):
        line = wrapped_lines[i]
        if i < len(wrapped_lines) - 1:  # Not the last line
            words_in_line = line.split()
            total_spaces = width - sum(len(word) for word in words_in_line)
            if len(words_in_line) > 1:
                space_between_words = total_spaces // (len(words_in_line) - 1)
                extra_spaces = total_spaces % (len(words_in_line) - 1)

                # Create justified line
                justified_line = ""
                for j in range(len(words_in_line)):
                    justified_line += words_in_line[j]
                    if j < len(words_in_line) - 1:  # Don't add space after the last word
                        justified_line += " " * (space_between_words + (1 if j < extra_spaces else 0))
                justified_lines.append(justified_line)
            else:
                justified_lines.append(line)  # Just a single word
        else:
            justified_lines.append(line)  # The last line remains left-aligned

    return justified_lines

# Upload file
uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPTX)", type=["pdf", "docx", "pptx"])
if uploaded_file is not None:
    # Save the uploaded file temporarily
    file_extension = uploaded_file.name.split('.')[-1].lower()
    file_path = f"temp_file.{file_extension}"
    with open(file_path, "wb") as f:
        f.write(uploaded_file.read())
    st.success("File uploaded successfully!")

    # Radio buttons for summary type selection
    summary_type = st.radio("Select summary type:", 
                             ("Number of Lines", "Number of Paragraphs", "Custom Prompt"))

    # Input fields based on summary type
    if summary_type == "Number of Lines":
        num_lines = st.number_input("Number of lines for summary:", min_value=1, max_value=20, value=5)

    elif summary_type == "Number of Paragraphs":
        num_paragraphs = st.number_input("Number of paragraphs for summary:", min_value=1, max_value=10, value=1)

    elif summary_type == "Custom Prompt":
        custom_prompt = st.text_input("Enter your custom prompt:", "Summarize the file content.")

    # Create a GenerativeModel instance
    model = genai.GenerativeModel("gemini-1.5-flash")

    # Button to generate summary
    if st.button("Get Summary"):
        try:
            # Extract text based on file type
            if file_extension == "pdf":
                file_text = extract_text_from_pdf(file_path)
            elif file_extension == "docx":
                file_text = extract_text_from_docx(file_path)
            elif file_extension == "pptx":
                file_text = extract_text_from_pptx(file_path)
            else:
                st.error("Unsupported file type.")
                raise ValueError("Unsupported file type.")

            # Prepare the prompt based on the selected summary type
            if summary_type == "Number of Lines":
                prompt = f"Give me a summary of this file in {num_lines} lines."
            elif summary_type == "Number of Paragraphs":
                prompt = f"Give me a summary of this file in {num_paragraphs} paragraphs."
            else:  # Custom Prompt
                prompt = custom_prompt

            # Generate content (summary) using the Gemini API
            response = model.generate_content([prompt, file_text])

            # Display the summary
            st.subheader("Summary Result:")
            if hasattr(response, 'text') and response.text:
                st.markdown(response.text)  # Show the summary in markdown for better readability
                
                # Create a download button for the summary
                buffer = io.BytesIO()
                pdf_filename = f"{uploaded_file.name.split('.')[0]}_summary.pdf"  # Set download file name

                # Use fitz to create a new PDF for the summary
                pdf_document = fitz.open()
                pdf_page = pdf_document.new_page()

                # Set up formatting for the PDF
                formatted_text = f"Summary of {uploaded_file.name}\n\n" + response.text
                
                # Replace double line breaks with a unique string to indicate paragraph breaks
                formatted_text = formatted_text.replace("\n\n", "\n\n[PARAGRAPH_BREAK]\n\n")

                # Wrap and justify text based on desired width
                wrapped_lines = wrap_and_justify_text(formatted_text.replace("[PARAGRAPH_BREAK]", "\n\n"), width=80)  # Replace marker with a new line for wrapping
                
                # Insert each wrapped line into the PDF
                y_position = 50  # Starting y position
                for line in wrapped_lines:
                    # If the line indicates a paragraph break, insert an empty line
                    if line.strip() == "":
                        y_position += 15  # Increase space for the paragraph break (empty line)
                    else:
                        pdf_page.insert_text((50, y_position), line, fontsize=12, fontname="helv", color=(0, 0, 0))
                        y_position += 15  # Move down for the next line (adjust as needed)

                    # Add an additional empty line after each line to create a gap
                    y_position += 5  # Additional space between lines

                # Save the PDF to a BytesIO buffer
                pdf_document.save(buffer)
                pdf_document.close()
                buffer.seek(0)

                # Download button
                st.download_button(
                    label="Download Summary as PDF",
                    data=buffer,
                    file_name=pdf_filename,
                    mime="application/pdf"
                )
            else:
                st.error("No summary returned. Please check the file content or try another document.")

        except Exception as e:
            st.error(f"An error occurred while generating the summary: {str(e)}")
